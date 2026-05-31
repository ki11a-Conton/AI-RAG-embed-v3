"""tests/test_doc_loader.py"""
import os
from lib.doc_loader import split_text_by_paragraphs, _extract_heading_path


def test_split_text_by_paragraphs_not_empty():
    text = "第一段内容。\n\n第二段内容。\n\n第三段内容。"
    chunks = split_text_by_paragraphs(text, chunk_size=20, chunk_overlap=5)
    assert len(chunks) > 0
    assert all(isinstance(c, tuple) and len(c) == 2 for c in chunks)
    assert all(c[0].strip() for c in chunks)


def test_split_text_basic():
    text = "Short para.\n\nAnother para.\n\nThird para."
    chunks = split_text_by_paragraphs(text, chunk_size=500, chunk_overlap=50)
    # 段落总长度小于 chunk_size，应合并为 1 个
    assert len(chunks) == 1
    assert "Short para." in chunks[0][0]


def test_split_text_large_paragraph():
    """段落超过 chunk_size 时应被进一步切分"""
    long_para = "A" * 2000
    chunks = split_text_by_paragraphs(long_para, chunk_size=500, chunk_overlap=50)
    assert len(chunks) > 1
    for chunk_text, _ in chunks:
        assert len(chunk_text) <= 500


def test_split_text_overlap():
    """最后一段应被保留"""
    text = "Para one.\n\nPara two.\n\nPara three."
    chunks = split_text_by_paragraphs(text, chunk_size=30, chunk_overlap=5)
    joined = " ".join(c for c, _ in chunks)
    assert "Para three" in joined


def test_extract_heading_path_with_headings():
    md_text = "# 自动控制原理\n## 二阶系统\n### 时间常数\n内容..."
    path = _extract_heading_path(md_text)
    assert path is not None
    assert "自动控制原理" in path
    assert ">" in path


def test_extract_heading_path_no_headings():
    text = "这是一段没有标题的正文内容。"
    path = _extract_heading_path(text)
    assert path is None


# ──────────────────────────────────────────────
# CSV 测试
# ──────────────────────────────────────────────

def test_read_csv_key_value_format():
    """Bug 回归：CSV 应输出 key: value 格式，header 行不作为数据行"""
    import tempfile, os
    from lib.doc_loader import _read_csv

    content = "name,age,role\nAlice,30,Engineer\nBob,25,Manager\n"
    with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False, encoding="utf-8") as f:
        f.write(content)
        path = f.name

    try:
        result = _read_csv(path)
        assert len(result) == 1
        text = result[0]["text"]
        # header 行不应单独出现
        assert "name,age,role" not in text
        assert "name: Alice" in text or "name: Alice" in text
        # 应包含字段名
        assert "name:" in text
        assert "age:" in text
        assert "role:" in text
        # 应包含值
        assert "Alice" in text
        assert "Engineer" in text
    finally:
        os.unlink(path)


def test_read_csv_empty():
    """空 CSV 文件应返回空列表"""
    import tempfile, os
    from lib.doc_loader import _read_csv

    with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False, encoding="utf-8") as f:
        f.write("name,age\n")  # 只有 header，无数据行
        path = f.name

    try:
        result = _read_csv(path)
        assert result == []
    finally:
        os.unlink(path)


# ──────────────────────────────────────────────
# CSV 按行分块测试 (T-10)
# ──────────────────────────────────────────────


def test_read_csv_small_with_chunk_size_still_one_page():
    """小 CSV + chunk_size：行为应与无 chunk_size 一致，返回单 page。"""
    import tempfile, os
    from lib.doc_loader import _read_csv

    content = "name,age\nAlice,30\nBob,25\n"
    with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False, encoding="utf-8") as f:
        f.write(content)
        path = f.name

    try:
        result = _read_csv(path, chunk_size=500)
        assert len(result) == 1
        assert result[0]["page"] is None
        text = result[0]["text"]
        assert "name: Alice" in text
        assert "age: 30" in text
        assert "name: Bob" in text
        assert "age: 25" in text
        # header 行不应出现在输出中
        assert "name,age" not in text
    finally:
        os.unlink(path)


def test_read_csv_large_returns_multiple_pages():
    """大 CSV + 小 chunk_size 应返回多个 page dict。"""
    import tempfile, os
    from lib.doc_loader import _read_csv

    # 生成 30 行，每行约 80 字符
    lines = ["col_a,col_b"]
    for i in range(30):
        lines.append(f"value_a_{i:03d},This is row number {i:03d} with padding text so it's not tiny")
    content = "\n".join(lines)

    with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False, encoding="utf-8") as f:
        f.write(content)
        path = f.name

    try:
        result = _read_csv(path, chunk_size=300)
        assert len(result) > 1, f"Expected multiple pages, got {len(result)}"
        for page in result:
            assert isinstance(page, dict)
            assert "text" in page
            assert page["page"] is None
            assert len(page["text"]) > 0
    finally:
        os.unlink(path)


def test_read_csv_chunk_never_splits_row():
    """任一 page 内不应出现被截断的行：每个 \n\n 分隔的段落都是完整行。"""
    import tempfile, os
    from lib.doc_loader import _read_csv

    # 生成 50 行，每行有唯一标识
    lines = ["id,label"]
    for i in range(50):
        lines.append(f"row_{i:04d},Label-for-row-{i:04d}-with-some-extra-padding-to-make-it-longer")
    content = "\n".join(lines)

    with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False, encoding="utf-8") as f:
        f.write(content)
        path = f.name

    try:
        result = _read_csv(path, chunk_size=400)
        assert len(result) > 1

        all_fragments: list[str] = []
        for page in result:
            # 每个 page 的 text 由完整行用 \n\n 拼接
            paragraphs = [p for p in page["text"].split("\n\n") if p.strip()]
            for p in paragraphs:
                assert p.startswith("id: row_"), f"Unexpected row fragment: {p[:60]}..."
                all_fragments.append(p)

        # 所有 50 行都应出现（无遗漏）
        assert len(all_fragments) == 50
        for i in range(50):
            expected_prefix = f"id: row_{i:04d}"
            found = any(f.startswith(expected_prefix) for f in all_fragments)
            assert found, f"Missing row: {expected_prefix}"
    finally:
        os.unlink(path)


def test_read_csv_chunk_headers_preserved():
    """header 标签应在文本中以 key: value 形式出现。"""
    import tempfile, os
    from lib.doc_loader import _read_csv

    content = "product,price,quantity\nApple,1.50,100\nBanana,0.75,200\n"
    with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False, encoding="utf-8") as f:
        f.write(content)
        path = f.name

    try:
        result = _read_csv(path, chunk_size=500)
        text = result[0]["text"]
        assert "product:" in text
        assert "price:" in text
        assert "quantity:" in text
        assert "product: Apple" in text
        assert "price: 1.50" in text
        assert "quantity: 100" in text
        # 原始 CSV header 行不应出现
        assert "product,price,quantity" not in text
    finally:
        os.unlink(path)


def test_read_csv_chunk_size_none_unchanged():
    """chunk_size=None 时保持原有单 page 行为不变。"""
    import tempfile, os
    from lib.doc_loader import _read_csv

    content = "col1,col2\na,b\nc,d\ne,f\n"
    with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False, encoding="utf-8") as f:
        f.write(content)
        path = f.name

    try:
        result = _read_csv(path, chunk_size=None)
        assert len(result) == 1
        assert result[0]["page"] is None
    finally:
        os.unlink(path)


def test_read_csv_oversized_row_does_not_crash():
    """单行超过 chunk_size 时不应崩溃，该行独占一个 page。"""
    import tempfile, os
    from lib.doc_loader import _read_csv

    # 一行极长，一行正常
    long_value = "X" * 600
    content = f"col1,col2\n{10},{long_value}\nnormal,row\n"
    with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False, encoding="utf-8") as f:
        f.write(content)
        path = f.name

    try:
        result = _read_csv(path, chunk_size=300)
        # 至少 2 个 page（长行独占一个，正常行另一个）
        assert len(result) >= 2
        # 长行应完整出现在某个 page 中
        all_text = "\n".join(p["text"] for p in result)
        assert long_value in all_text
        assert "normal" in all_text
    finally:
        os.unlink(path)


def test_read_txt_md_invalid_utf8():
    """_read_txt_md should not raise on malformed UTF-8 bytes."""
    import tempfile, os
    from lib.doc_loader import _read_txt_md

    # Invalid UTF-8 bytes mixed with readable ASCII.
    raw = b"\xff\xfeHello\xffworld\n\x80\x81line two\n"
    with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as f:
        f.write(raw)
        path = f.name

    try:
        result = _read_txt_md(path)
        assert len(result) == 1
        text = result[0]["text"]
        assert "Hello" in text or "line" in text
    finally:
        os.unlink(path)


def test_read_txt_md_invalid_utf8_md():
    """read_file_pages should apply the same tolerant .md handling."""
    import tempfile, os
    from lib.doc_loader import read_file_pages

    raw = b"\xff\xfe# Title\xff\nContent with \x80invalid bytes.\n"
    with tempfile.NamedTemporaryFile(suffix=".md", delete=False) as f:
        f.write(raw)
        path = f.name

    try:
        result = read_file_pages(path)
        assert len(result) == 1
        text = result[0]["text"]
        assert "Title" in text or "Content" in text
    finally:
        os.unlink(path)


def test_split_markdown_heading_stack():
    """Bug 回归：同级标题切换时 heading_stack 应正确弹出"""
    from lib.doc_loader import _split_markdown_by_headings

    md = "# Part A\n## Section 1\n### Sub 1.1\ncontent 1.1\n## Section 2\ncontent 2\n"
    result = _split_markdown_by_headings(md, chunk_size=500, chunk_overlap=50)
    paths = [p for p, _, _ in result]
    # Section 2 不应携带 Sub 1.1 的路径
    section2_path = next((p for p, t, _ in result if "content 2" in t), None)
    assert section2_path is not None
    assert "Sub 1.1" not in section2_path
    assert "Section 2" in section2_path


def test_load_documents_show_progress_emits_output(capsys, tmp_path):
    """show_progress should emit progress while the default stays quiet."""
    from lib.doc_loader import load_documents

    doc1 = tmp_path / "doc1.txt"
    doc1.write_text(
        "Hello world.\n\nSecond paragraph.\n\nThird paragraph.",
        encoding="utf-8",
    )
    doc2 = tmp_path / "doc2.txt"
    doc2.write_text("Fourth paragraph.\n\nFifth paragraph.", encoding="utf-8")

    chunks_default = load_documents(str(tmp_path), chunk_size=200, chunk_overlap=20)
    captured_default = capsys.readouterr().out
    assert "chunks from" not in captured_default
    assert len(chunks_default) > 0

    chunks_progress = load_documents(
        str(tmp_path),
        chunk_size=200,
        chunk_overlap=20,
        show_progress=True,
    )
    captured_progress = capsys.readouterr().out
    assert "chunks from" in captured_progress
    assert "files" in captured_progress
    assert len(chunks_progress) > 0


# ──────────────────────────────────────────────
# _table_to_text 测试
# ──────────────────────────────────────────────


def test_table_to_text_basic():
    """标准表格行应格式化为 header: value 样式。"""
    from lib.doc_loader import _table_to_text

    headers = ["Name", "Age", "Role"]
    rows = [["Alice", "30", "Engineer"], ["Bob", "25", "Manager"]]
    result = _table_to_text(headers, rows)

    assert "Name: Alice" in result
    assert "Age: 30" in result
    assert "Role: Engineer" in result
    assert "Name: Bob" in result
    assert "Age: 25" in result
    assert "Role: Manager" in result
    # 两行之间应有换行
    assert "\n" in result


def test_table_to_text_skips_empty_values():
    """空值不应出现在输出中。"""
    from lib.doc_loader import _table_to_text

    headers = ["Field A", "Field B", "Field C"]
    rows = [["hello", "", ""], ["", "world", ""], ["", "", ""]]
    result = _table_to_text(headers, rows)

    lines = [l for l in result.split("\n") if l.strip()]
    assert len(lines) == 2
    # 第一行只有 Field A，不含 B 和 C
    assert "Field A: hello" in lines[0]
    assert "Field B" not in lines[0]
    assert "Field C" not in lines[0]
    # 第二行只有 Field B，不含 A 和 C
    assert "Field B: world" in lines[1]
    assert "Field A" not in lines[1]
    assert "Field C" not in lines[1]


def test_table_to_text_all_empty_rows():
    """全空行应返回空字符串。"""
    from lib.doc_loader import _table_to_text

    headers = ["X", "Y"]
    rows = [["", ""], [None, ""]]
    result = _table_to_text(headers, rows)
    assert result == ""


def test_table_to_text_no_rows():
    """无数据行应返回空字符串。"""
    from lib.doc_loader import _table_to_text

    headers = ["X", "Y"]
    result = _table_to_text(headers, [])
    assert result == ""


# ──────────────────────────────────────────────
# DOCX 表格提取测试
# ──────────────────────────────────────────────


def test_read_docx_with_tables():
    """_read_docx 应在段落之外同时提取表格内容。"""
    import tempfile
    import os
    from docx import Document
    from lib.doc_loader import _read_docx

    doc = Document()
    doc.add_paragraph("First paragraph.")
    doc.add_paragraph("Second paragraph.")

    table = doc.add_table(rows=3, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Age"
    table.cell(1, 0).text = "Alice"
    table.cell(1, 1).text = "30"
    table.cell(2, 0).text = "Bob"
    table.cell(2, 1).text = "25"

    path = os.path.join(tempfile.gettempdir(), "test_table_doc_loader.docx")
    doc.save(path)

    try:
        result = _read_docx(path)
        assert len(result) == 1
        text = result[0]["text"]
        assert "First paragraph" in text
        assert "Second paragraph" in text
        assert "Name: Alice" in text
        assert "Age: 30" in text
        assert "Name: Bob" in text
        assert "Age: 25" in text
        assert result[0]["page"] is None
    finally:
        os.unlink(path)


def test_read_docx_table_with_empty_cells():
    """表格中空单元格应被跳过，不出现无值的 header: 片段。"""
    import tempfile
    import os
    from docx import Document
    from lib.doc_loader import _read_docx

    doc = Document()
    doc.add_paragraph("Doc with sparse table.")

    table = doc.add_table(rows=2, cols=3)
    table.cell(0, 0).text = "Col1"
    table.cell(0, 1).text = "Col2"
    table.cell(0, 2).text = "Col3"
    table.cell(1, 0).text = "val1"
    table.cell(1, 1).text = ""  # empty
    table.cell(1, 2).text = "val3"

    path = os.path.join(tempfile.gettempdir(), "test_sparse_table.docx")
    doc.save(path)

    try:
        result = _read_docx(path)
        text = result[0]["text"]
        assert "Col1: val1" in text
        assert "Col3: val3" in text
        assert "Col2:" not in text  # empty cell skipped
    finally:
        os.unlink(path)


def test_read_docx_no_tables():
    """纯段落 DOCX（无表格）行为不变。"""
    import tempfile
    import os
    from docx import Document
    from lib.doc_loader import _read_docx

    doc = Document()
    doc.add_paragraph("Only a paragraph here.")
    doc.add_paragraph("Another one.")

    path = os.path.join(tempfile.gettempdir(), "test_no_table.docx")
    doc.save(path)

    try:
        result = _read_docx(path)
        assert len(result) == 1
        text = result[0]["text"]
        assert "Only a paragraph here" in text
        assert "Another one" in text
    finally:
        os.unlink(path)


# ──────────────────────────────────────────────
# PPTX 测试
# ──────────────────────────────────────────────


def test_read_pptx_basic():
    """_read_pptx 应从两页幻灯片提取文本，page 为 1-based。"""
    import tempfile
    import os
    from pptx import Presentation
    from lib.doc_loader import _read_pptx

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]

    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Slide One Title"
    txBox = slide1.shapes.add_textbox(10, 10, 300, 50)
    txBox.text_frame.text = "Body text for slide one."

    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Slide Two Title"
    txBox2 = slide2.shapes.add_textbox(10, 10, 300, 50)
    txBox2.text_frame.text = "Body text for slide two."

    path = os.path.join(tempfile.gettempdir(), "test_basic_pptx.pptx")
    prs.save(path)

    try:
        result = _read_pptx(path)
        assert len(result) == 2
        assert result[0]["page"] == 1
        assert result[1]["page"] == 2
        assert "Slide One Title" in result[0]["text"]
        assert "Body text for slide one." in result[0]["text"]
        assert "Slide Two Title" in result[1]["text"]
        assert "Body text for slide two." in result[1]["text"]
    finally:
        os.unlink(path)


def test_read_pptx_empty_slides_skipped():
    """空白幻灯片应被跳过。"""
    import tempfile
    import os
    from pptx import Presentation
    from lib.doc_loader import _read_pptx

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]

    # empty slide (no text)
    prs.slides.add_slide(slide_layout)

    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Only content slide"

    path = os.path.join(tempfile.gettempdir(), "test_empty_pptx.pptx")
    prs.save(path)

    try:
        result = _read_pptx(path)
        assert len(result) == 1
        assert result[0]["page"] == 2
        assert "Only content slide" in result[0]["text"]
    finally:
        os.unlink(path)


def test_read_pptx_with_table():
    """_read_pptx 应提取幻灯片中的表格文本。"""
    import tempfile
    import os
    from pptx import Presentation
    from pptx.util import Inches
    from lib.doc_loader import _read_pptx

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Slide with Table"

    rows, cols = 3, 2
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(4), Inches(3))
    table = table_shape.table
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Score"
    table.cell(1, 0).text = "Alice"
    table.cell(1, 1).text = "95"
    table.cell(2, 0).text = "Bob"
    table.cell(2, 1).text = "87"

    path = os.path.join(tempfile.gettempdir(), "test_table_pptx.pptx")
    prs.save(path)

    try:
        result = _read_pptx(path)
        assert len(result) == 1
        assert result[0]["page"] == 1
        text = result[0]["text"]
        assert "Slide with Table" in text
        assert "Name: Alice" in text
        assert "Score: 95" in text
        assert "Name: Bob" in text
        assert "Score: 87" in text
    finally:
        os.unlink(path)


def test_read_pptx_import_error():
    """python-pptx 未安装时应抛出友好的 ImportError。"""
    import builtins
    from unittest.mock import patch
    from lib.doc_loader import _read_pptx

    real_import = builtins.__import__

    def fake_import(name, *args, **kwargs):
        if name == "pptx":
            raise ImportError("No module named pptx")
        return real_import(name, *args, **kwargs)

    with patch.object(builtins, "__import__", side_effect=fake_import):
        try:
            _read_pptx("fake.pptx")
            assert False, "Should have raised ImportError"
        except ImportError as e:
            assert "python-pptx" in str(e)


def test_read_pptx_via_read_file_pages():
    """read_file_pages 应正确路由 .pptx 文件到 _read_pptx。"""
    import tempfile
    import os
    from pptx import Presentation
    from lib.doc_loader import read_file_pages

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Routed via read_file_pages"

    path = os.path.join(tempfile.gettempdir(), "test_routing_pptx.pptx")
    prs.save(path)

    try:
        result = read_file_pages(path)
        assert len(result) == 1
        assert result[0]["page"] == 1
        assert "Routed via read_file_pages" in result[0]["text"]
    finally:
        os.unlink(path)


def test_read_pptx_page_never_none():
    """PPTX 结果的 page 字段应始终为 1-based int，不为 None。"""
    import tempfile
    import os
    from pptx import Presentation
    from lib.doc_loader import _read_pptx

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Content"

    path = os.path.join(tempfile.gettempdir(), "test_page_none_pptx.pptx")
    prs.save(path)

    try:
        result = _read_pptx(path)
        assert len(result) == 1
        assert isinstance(result[0]["page"], int)
        assert result[0]["page"] >= 1
    finally:
        os.unlink(path)


# ──────────────────────────────────────────────
# PDF 测试
# ──────────────────────────────────────────────


def test_read_pdf_prefers_pymupdf():
    """PyMuPDF 可用时应优先使用，且 page 应为 1-based，空页跳过。"""
    from unittest.mock import patch, MagicMock
    import sys
    from lib.doc_loader import _read_pdf

    # 构造 mock fitz 模块
    mock_fitz = MagicMock()
    mock_page1 = MagicMock()
    mock_page1.get_text.return_value = "Page one content."
    mock_page2 = MagicMock()
    mock_page2.get_text.return_value = "   "  # only whitespace → empty
    mock_page3 = MagicMock()
    mock_page3.get_text.return_value = "Page three content."

    mock_doc = MagicMock()
    mock_doc.__iter__.return_value = iter([mock_page1, mock_page2, mock_page3])
    mock_fitz.open.return_value = mock_doc

    with patch.dict(sys.modules, {"fitz": mock_fitz}):
        result = _read_pdf("fake_path.pdf")

    assert len(result) == 2  # page 2 (empty) skipped
    assert result[0]["page"] == 1
    assert result[1]["page"] == 3
    assert result[0]["text"] == "Page one content."
    assert result[1]["text"] == "Page three content."
    mock_doc.close.assert_called_once()


def test_read_pdf_fallback_to_pypdf():
    """PyMuPDF 未安装时应回退到 pypdf，保持 1-based page。"""
    from unittest.mock import patch, MagicMock
    import builtins
    import sys
    import types
    from lib.doc_loader import _read_pdf

    real_import = builtins.__import__
    mock_reader = MagicMock()
    mock_page = MagicMock()
    mock_page.extract_text.return_value = "PDF content."
    mock_reader.pages = [mock_page]

    mock_pypdf = types.SimpleNamespace(PdfReader=MagicMock(return_value=mock_reader))

    def fake_import(name, *args, **kwargs):
        if name == "fitz":
            raise ImportError("fitz unavailable")
        if name == "pypdf":
            return mock_pypdf
        return real_import(name, *args, **kwargs)

    with patch.object(builtins, "__import__", side_effect=fake_import):
        result = _read_pdf("fake_path.pdf")

    assert len(result) == 1
    assert result[0]["page"] == 1
    assert result[0]["text"] == "PDF content."
    mock_page.extract_text.assert_called_once_with(extraction_mode="layout")
    mock_pypdf.PdfReader.assert_called_once_with("fake_path.pdf")


def test_read_pdf_skips_empty_pages_pypdf():
    """pypdf 回退路径同样应跳过空白页。"""
    from unittest.mock import patch, MagicMock
    import builtins
    import types
    from lib.doc_loader import _read_pdf

    real_import = builtins.__import__
    mock_reader = MagicMock()
    mock_page1 = MagicMock()
    mock_page1.extract_text.return_value = ""
    mock_page2 = MagicMock()
    mock_page2.extract_text.return_value = "Non-empty."
    mock_page3 = MagicMock()
    mock_page3.extract_text.return_value = "   \n  "
    mock_reader.pages = [mock_page1, mock_page2, mock_page3]

    mock_pypdf = types.SimpleNamespace(PdfReader=MagicMock(return_value=mock_reader))

    def fake_import(name, *args, **kwargs):
        if name == "fitz":
            raise ImportError("fitz unavailable")
        if name == "pypdf":
            return mock_pypdf
        return real_import(name, *args, **kwargs)

    with patch.object(builtins, "__import__", side_effect=fake_import):
        result = _read_pdf("fake_path.pdf")

    assert len(result) == 1
    assert result[0]["page"] == 2
    assert result[0]["text"] == "Non-empty."


def test_read_pdf_page_none_never():
    """PDF 结果的 page 字段应始终为 1-based int，不为 None。"""
    from unittest.mock import patch, MagicMock
    import sys
    from lib.doc_loader import _read_pdf

    mock_fitz = MagicMock()
    mock_page = MagicMock()
    mock_page.get_text.return_value = "Content"
    mock_doc = MagicMock()
    mock_doc.__iter__.return_value = iter([mock_page])
    mock_fitz.open.return_value = mock_doc

    with patch.dict(sys.modules, {"fitz": mock_fitz}):
        result = _read_pdf("fake_path.pdf")

    assert len(result) == 1
    assert isinstance(result[0]["page"], int)
    assert result[0]["page"] >= 1


# ──────────────────────────────────────────────
# 图片提取（image caption）集成测试
# ──────────────────────────────────────────────


def test_load_documents_disabled_default_no_caption_chunks(tmp_path):
    """Default (image_captioner=None) must produce zero image_caption chunks."""
    from lib.doc_loader import load_documents

    txt = tmp_path / "readme.txt"
    txt.write_text("Hello world.", encoding="utf-8")

    chunks = load_documents(str(tmp_path), chunk_size=200, chunk_overlap=20)
    file_types = {c["metadata"]["file_type"] for c in chunks}
    assert "image_caption" not in file_types


def test_load_documents_with_mocked_image_captioner_adds_caption_chunks():
    """When an image_captioner is supplied, PDF images produce caption chunks."""
    from unittest.mock import MagicMock, patch
    from lib.doc_loader import load_documents

    mock_cap = MagicMock()
    mock_cap.caption.return_value = "A flowchart diagram."

    fake_pdf_images = [(b"img1bytes", 1), (b"img2bytes", 3)]

    # Patch _read_pdf so the file is not skipped, and _extract_pdf_images
    # so we control what images are returned.
    with patch(
        "lib.doc_loader._read_pdf",
        return_value=[{"text": "Page text.", "page": 1}],
    ), patch(
        "lib.doc_loader._extract_pdf_images",
        return_value=fake_pdf_images,
    ):
        # Create a dummy .pdf file so the extension check passes
        import tempfile, os
        with tempfile.TemporaryDirectory() as tmpdir:
            pdf = os.path.join(tmpdir, "sample.pdf")
            with open(pdf, "wb") as f:
                f.write(b"%PDF-fake")

            chunks = load_documents(
                tmpdir, chunk_size=500, chunk_overlap=50, image_captioner=mock_cap
            )

    # At least 2 caption chunks (plus the text chunk from the PDF page)
    caption_chunks = [c for c in chunks if c["metadata"]["file_type"] == "image_caption"]
    assert len(caption_chunks) == 2

    c0 = caption_chunks[0]
    assert c0["metadata"]["source"] == "sample.pdf"
    assert c0["metadata"]["page"] == 1
    assert c0["metadata"]["image_index"] == 0
    assert c0["text"] == "A flowchart diagram."

    c1 = caption_chunks[1]
    assert c1["metadata"]["page"] == 3
    assert c1["metadata"]["image_index"] == 1

    assert mock_cap.caption.call_count == 2


def test_image_caption_chunk_metadata():
    """Caption chunks must carry the required metadata fields."""
    from unittest.mock import MagicMock, patch
    from lib.doc_loader import load_documents
    import tempfile, os

    mock_cap = MagicMock()
    mock_cap.caption.return_value = "Extracted image description."

    with patch(
        "lib.doc_loader._read_pdf",
        return_value=[{"text": "page text", "page": 2}],
    ), patch(
        "lib.doc_loader._extract_pdf_images",
        return_value=[(b"img", 2)],
    ):
        with tempfile.TemporaryDirectory() as tmpdir:
            pdf = os.path.join(tmpdir, "doc.pdf")
            with open(pdf, "wb") as f:
                f.write(b"%PDF-fake")

            chunks = load_documents(
                tmpdir, chunk_size=500, chunk_overlap=50, image_captioner=mock_cap
            )

    caption_chunks = [c for c in chunks if c["metadata"]["file_type"] == "image_caption"]
    assert len(caption_chunks) >= 1
    meta = caption_chunks[0]["metadata"]
    assert meta["file_type"] == "image_caption"
    assert meta["source"] == "doc.pdf"
    assert meta["page"] == 2
    assert meta["image_index"] == 0


def test_image_captioner_not_called_for_non_pdf_docx():
    """Image captioner must never be invoked for .txt, .md, etc."""
    from unittest.mock import MagicMock
    from lib.doc_loader import load_documents
    import tempfile, os

    mock_cap = MagicMock()

    with tempfile.TemporaryDirectory() as tmpdir:
        txt = os.path.join(tmpdir, "notes.txt")
        with open(txt, "w", encoding="utf-8") as f:
            f.write("Plain text content.")
        md = os.path.join(tmpdir, "notes.md")
        with open(md, "w", encoding="utf-8") as f:
            f.write("# Heading\n\nBody.")

        chunks = load_documents(
            tmpdir, chunk_size=200, chunk_overlap=20, image_captioner=mock_cap
        )

    mock_cap.caption.assert_not_called()
    assert all(c["metadata"]["file_type"] != "image_caption" for c in chunks)


def test_image_caption_context_hint_includes_page():
    """The context_hint passed to caption() should mention page when known."""
    from unittest.mock import MagicMock, patch
    from lib.doc_loader import load_documents
    import tempfile, os

    mock_cap = MagicMock()
    mock_cap.caption.return_value = "ok"

    with patch(
        "lib.doc_loader._read_pdf",
        return_value=[{"text": "text", "page": 4}],
    ), patch(
        "lib.doc_loader._extract_pdf_images",
        return_value=[(b"img", 4)],
    ):
        with tempfile.TemporaryDirectory() as tmpdir:
            pdf = os.path.join(tmpdir, "file.pdf")
            with open(pdf, "wb") as f:
                f.write(b"%PDF-fake")

            load_documents(
                tmpdir, chunk_size=500, chunk_overlap=50, image_captioner=mock_cap
            )

    mock_cap.caption.assert_called_once()
    # context_hint should mention page 4
    hint = mock_cap.caption.call_args.kwargs.get("context_hint", "")
    assert "page 4" in hint
    assert "file.pdf" in hint


def test_extract_pdf_images_returns_expected_structure():
    """Unit test for the PDF image extraction helper."""
    from unittest.mock import patch, MagicMock
    import sys
    from lib.doc_loader import _extract_pdf_images

    mock_page1 = MagicMock()
    mock_page1.get_images.return_value = [(5, 0, 0, 0, 0, 0, 0)]

    mock_page2 = MagicMock()
    mock_page2.get_images.return_value = []  # no images on page 2

    mock_page3 = MagicMock()
    mock_page3.get_images.return_value = [
        (7, 0, 0, 0, 0, 0, 0),
        (8, 0, 0, 0, 0, 0, 0),
    ]

    mock_doc = MagicMock()
    mock_doc.__iter__.return_value = iter([mock_page1, mock_page2, mock_page3])
    mock_doc.extract_image.side_effect = lambda xref: {
        5: {"image": b"image_on_page_1"},
        7: {"image": b"first_image_page_3"},
        8: {"image": b"second_image_page_3"},
    }[xref]

    mock_fitz = MagicMock()
    mock_fitz.open.return_value = mock_doc

    with patch.dict(sys.modules, {"fitz": mock_fitz}):
        results = _extract_pdf_images("fake.pdf")

    assert len(results) == 3
    assert results[0] == (b"image_on_page_1", 1)
    assert results[1] == (b"first_image_page_3", 3)
    assert results[2] == (b"second_image_page_3", 3)
    mock_doc.close.assert_called_once()


def test_extract_pdf_images_no_fitz_returns_empty():
    """When PyMuPDF is not installed, return an empty list."""
    from unittest.mock import patch
    import builtins
    from lib.doc_loader import _extract_pdf_images

    real_import = builtins.__import__

    def fake_import(name, *args, **kwargs):
        if name == "fitz":
            raise ImportError("no fitz")
        return real_import(name, *args, **kwargs)

    with patch.object(builtins, "__import__", side_effect=fake_import):
        results = _extract_pdf_images("any.pdf")

    assert results == []


# ──────────────────────────────────────────────
# Markdown heading prefix language config (T-13)
# ──────────────────────────────────────────────


def test_load_documents_md_default_prefix_is_chinese(tmp_path):
    """Default (docs_lang='zh') uses Chinese heading prefix '标题路径：'."""
    from lib.doc_loader import load_documents

    md = tmp_path / "doc.md"
    md.write_text("# Chapter 1\n\nSome content here.\n\n## Section A\n\nMore details.\n", encoding="utf-8")

    chunks = load_documents(str(tmp_path), chunk_size=500, chunk_overlap=50)

    # Locate chunks that came from a heading section
    heading_chunks = [c for c in chunks if c["text"].startswith("标题路径：")]
    assert len(heading_chunks) > 0, "Expected some chunks with Chinese heading prefix"

    for c in heading_chunks:
        assert "标题路径：" in c["text"]
        # Metadata must be untouched
        assert "source" in c["metadata"]
        assert c["metadata"]["file_type"] == "md"


def test_load_documents_md_english_prefix(tmp_path):
    """docs_lang='en' uses English heading prefix 'Heading: '."""
    from lib.doc_loader import load_documents

    md = tmp_path / "doc.md"
    md.write_text("# Chapter 1\n\nSome content here.\n\n## Section A\n\nMore details.\n", encoding="utf-8")

    chunks = load_documents(str(tmp_path), chunk_size=500, chunk_overlap=50, docs_lang="en")

    heading_chunks = [c for c in chunks if "Heading:" in c["text"]]
    assert len(heading_chunks) > 0, "Expected some chunks with English heading prefix"

    for c in heading_chunks:
        assert "Heading:" in c["text"]
        # Must NOT contain the Chinese prefix
        assert "标题路径" not in c["text"]
        # Metadata must be untouched
        assert "source" in c["metadata"]
        assert c["metadata"]["file_type"] == "md"


def test_load_documents_md_no_heading_chunk_no_prefix(tmp_path):
    """Chunks without a heading path must have no prefix regardless of docs_lang."""
    from lib.doc_loader import load_documents

    md = tmp_path / "doc.md"
    md.write_text("Plain text without any markdown headings.\n\nSecond paragraph.\n", encoding="utf-8")

    chunks_default = load_documents(str(tmp_path), chunk_size=500, chunk_overlap=50)
    chunks_en = load_documents(str(tmp_path), chunk_size=500, chunk_overlap=50, docs_lang="en")

    for chunks, label in [(chunks_default, "zh"), (chunks_en, "en")]:
        for c in chunks:
            assert not c["text"].startswith("标题路径："), f"{label}: unexpected Chinese prefix"
            assert not c["text"].startswith("Heading:"), f"{label}: unexpected English prefix"


def test_load_documents_md_heading_prefix_metadata_unchanged(tmp_path):
    """Metadata fields must be identical regardless of docs_lang setting."""
    from lib.doc_loader import load_documents

    md = tmp_path / "doc.md"
    md.write_text("# Title\n\nBody text.\n", encoding="utf-8")

    chunks_zh = load_documents(str(tmp_path), chunk_size=500, chunk_overlap=50, docs_lang="zh")
    chunks_en = load_documents(str(tmp_path), chunk_size=500, chunk_overlap=50, docs_lang="en")

    assert len(chunks_zh) == len(chunks_en)
    for cz, ce in zip(chunks_zh, chunks_en):
        assert cz["metadata"] == ce["metadata"], "Metadata must not change with docs_lang"
        # The text must differ (different prefix) but the content after the prefix must match
        assert cz["text"] != ce["text"]
        # Both should contain the heading title
        assert "Title" in cz["text"]
        assert "Title" in ce["text"]


# ──────────────────────────────────────────────
# XLSX 测试 (T-06)
# ──────────────────────────────────────────────


def test_read_xlsx_basic(tmp_path):
    """_read_xlsx should extract rows as key:value text per sheet."""
    import openpyxl
    from lib.doc_loader import _read_xlsx

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Age", "Role"])
    ws.append(["Alice", "30", "Engineer"])
    ws.append(["Bob", "25", "Manager"])

    xlsx_path = tmp_path / "test.xlsx"
    wb.save(str(xlsx_path))
    wb.close()

    result = _read_xlsx(str(xlsx_path))
    assert len(result) == 1
    text = result[0]["text"]
    assert "Sheet: Sheet1" in text
    assert "Name: Alice" in text
    assert "Age: 30" in text
    assert "Role: Engineer" in text
    assert "Name: Bob" in text
    assert "Age: 25" in text
    assert "Role: Manager" in text
    assert result[0]["page"] == 1


def test_read_xlsx_multiple_sheets(tmp_path):
    """Each sheet becomes a separate page with page=sheet_index."""
    import openpyxl
    from lib.doc_loader import _read_xlsx

    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Employees"
    ws1.append(["Name", "Dept"])
    ws1.append(["Alice", "Engineering"])

    ws2 = wb.create_sheet("Products")
    ws2.append(["Product", "Price"])
    ws2.append(["Widget", "9.99"])

    xlsx_path = tmp_path / "multi.xlsx"
    wb.save(str(xlsx_path))
    wb.close()

    result = _read_xlsx(str(xlsx_path))
    assert len(result) == 2
    assert result[0]["page"] == 1
    assert result[1]["page"] == 2
    assert "Sheet: Employees" in result[0]["text"]
    assert "Alice" in result[0]["text"]
    assert "Sheet: Products" in result[1]["text"]
    assert "Widget" in result[1]["text"]


def test_read_xlsx_empty_sheet_skipped(tmp_path):
    """Empty sheets should not produce pages."""
    import openpyxl
    from lib.doc_loader import _read_xlsx

    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "NonEmpty"
    ws1.append(["Col1"])
    ws1.append(["data"])

    ws2 = wb.create_sheet("EmptySheet")
    # No data rows

    xlsx_path = tmp_path / "empty.xlsx"
    wb.save(str(xlsx_path))
    wb.close()

    result = _read_xlsx(str(xlsx_path))
    assert len(result) == 1
    assert "NonEmpty" in result[0]["text"]
    assert result[0]["page"] == 1


def test_read_xlsx_import_error():
    """When openpyxl is not installed, raise a friendly ImportError."""
    import builtins
    from unittest.mock import patch
    from lib.doc_loader import _read_xlsx

    real_import = builtins.__import__

    def fake_import(name, *args, **kwargs):
        if name == "openpyxl":
            raise ImportError("No module named openpyxl")
        return real_import(name, *args, **kwargs)

    with patch.object(builtins, "__import__", side_effect=fake_import):
        try:
            _read_xlsx("fake.xlsx")
            assert False, "Should have raised ImportError"
        except ImportError as e:
            assert "openpyxl" in str(e)


# ──────────────────────────────────────────────
# 独立图片文件测试 (T-06)
# ──────────────────────────────────────────────


def test_image_file_without_captioner_no_chunks(tmp_path):
    """Image files without image_captioner should produce no chunks (no crash)."""
    from lib.doc_loader import load_documents
    from PIL import Image

    img = Image.new("RGB", (10, 10), color="red")
    img_path = tmp_path / "test.png"
    img.save(str(img_path))

    chunks = load_documents(str(tmp_path), chunk_size=200, chunk_overlap=20)
    # No chunks should be produced for a plain image
    assert len(chunks) == 0


def test_image_file_with_captioner_produces_caption_chunks(tmp_path):
    """With mocked captioner, image files produce caption chunks."""
    from unittest.mock import MagicMock
    from lib.doc_loader import load_documents
    from PIL import Image

    mock_cap = MagicMock()
    mock_cap.caption.return_value = "A red square image."

    img = Image.new("RGB", (10, 10), color="red")
    img_path = tmp_path / "test.png"
    img.save(str(img_path))

    chunks = load_documents(
        str(tmp_path), chunk_size=200, chunk_overlap=20, image_captioner=mock_cap
    )

    assert len(chunks) == 1
    assert chunks[0]["text"] == "A red square image."
    assert chunks[0]["metadata"]["file_type"] == "image_caption"
    assert chunks[0]["metadata"]["source"] == "test.png"
    assert chunks[0]["metadata"]["page"] is None
    assert chunks[0]["metadata"]["image_index"] == 0
    mock_cap.caption.assert_called_once()


def test_image_file_caption_context_hint(tmp_path):
    """The context_hint should include the filename."""
    from unittest.mock import MagicMock
    from lib.doc_loader import load_documents
    from PIL import Image

    mock_cap = MagicMock()
    mock_cap.caption.return_value = "desc"

    img = Image.new("RGB", (5, 5), color="blue")
    img_path = tmp_path / "photo.jpg"
    img.save(str(img_path))

    load_documents(
        str(tmp_path), chunk_size=200, chunk_overlap=20, image_captioner=mock_cap
    )

    mock_cap.caption.assert_called_once()
    hint = mock_cap.caption.call_args.kwargs.get("context_hint", "")
    assert "photo.jpg" in hint


def test_image_file_multiple_extensions(tmp_path):
    """All image extensions should be handled (png, jpg, jpeg, webp)."""
    from unittest.mock import MagicMock
    from lib.doc_loader import load_documents
    from PIL import Image

    mock_cap = MagicMock()
    mock_cap.caption.return_value = "desc"

    exts = [".png", ".jpg", ".jpeg", ".webp"]
    for ext in exts:
        img = Image.new("RGB", (3, 3), color="green")
        img_path = tmp_path / f"img{ext}"
        if ext == ".webp":
            img.save(str(img_path), format="WEBP")
        else:
            img.save(str(img_path))

    chunks = load_documents(
        str(tmp_path), chunk_size=200, chunk_overlap=20, image_captioner=mock_cap
    )

    caption_chunks = [c for c in chunks if c["metadata"]["file_type"] == "image_caption"]
    assert len(caption_chunks) == len(exts), (
        f"Expected {len(exts)} caption chunks, got {len(caption_chunks)}"
    )
    assert mock_cap.caption.call_count == len(exts)


def test_image_file_captioner_exception_handled(tmp_path):
    """When caption() raises, the image should be skipped without crashing."""
    from unittest.mock import MagicMock
    from lib.doc_loader import load_documents
    from PIL import Image

    mock_cap = MagicMock()
    mock_cap.caption.side_effect = RuntimeError("Caption failed")

    img = Image.new("RGB", (5, 5), color="yellow")
    img_path = tmp_path / "error.png"
    img.save(str(img_path))

    # Should not raise
    chunks = load_documents(
        str(tmp_path), chunk_size=200, chunk_overlap=20, image_captioner=mock_cap
    )

    # No caption chunks should be produced since caption failed
    caption_chunks = [c for c in chunks if c["metadata"]["file_type"] == "image_caption"]
    assert len(caption_chunks) == 0


# only_sources parameter tests


def test_load_documents_only_sources_skips_unselected(tmp_path):
    """only_sources filters out files whose relative path is not in the set."""
    from lib.doc_loader import load_documents

    keep1 = tmp_path / "keep.txt"
    keep1.write_text("This file should be included.", encoding="utf-8")

    keep2 = tmp_path / "sub" / "also_keep.md"
    keep2.parent.mkdir(parents=True, exist_ok=True)
    keep2.write_text("# Heading\n\nThis file should also be included.", encoding="utf-8")

    skip1 = tmp_path / "skip.txt"
    skip1.write_text("This file should be skipped.", encoding="utf-8")

    skip2 = tmp_path / "other" / "skip_too.md"
    skip2.parent.mkdir(parents=True, exist_ok=True)
    skip2.write_text("# Skipped\n\nThis file should be skipped.", encoding="utf-8")

    only_set = {"keep.txt", os.path.join("sub", "also_keep.md")}
    chunks = load_documents(
        str(tmp_path), chunk_size=500, chunk_overlap=50, only_sources=only_set,
    )

    sources = {c["metadata"]["source"] for c in chunks}
    assert "keep.txt" in sources, "keep.txt should be loaded"
    assert os.path.join("sub", "also_keep.md") in sources, "sub/also_keep.md should be loaded"
    assert "skip.txt" not in sources, "skip.txt should not be loaded"
    assert os.path.join("other", "skip_too.md") not in sources, "other/skip_too.md not loaded"


def test_load_documents_only_sources_no_filtering_when_none(tmp_path):
    """When only_sources is None, all supported files are loaded as before."""
    from lib.doc_loader import load_documents

    doc1 = tmp_path / "a.txt"
    doc1.write_text("File A.", encoding="utf-8")
    doc2 = tmp_path / "b.txt"
    doc2.write_text("File B.", encoding="utf-8")

    chunks = load_documents(str(tmp_path), chunk_size=200, chunk_overlap=20)
    sources = {c["metadata"]["source"] for c in chunks}
    assert "a.txt" in sources
    assert "b.txt" in sources


def test_load_documents_only_sources_empty_set_loads_nothing(tmp_path):
    """An empty only_sources set should produce zero chunks."""
    from lib.doc_loader import load_documents

    doc = tmp_path / "doc.txt"
    doc.write_text("Content.", encoding="utf-8")

    chunks = load_documents(
        str(tmp_path), chunk_size=200, chunk_overlap=20, only_sources=set(),
    )
    assert len(chunks) == 0
