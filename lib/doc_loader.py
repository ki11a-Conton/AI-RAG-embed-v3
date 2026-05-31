"""
doc_loader.py
支持 .txt / .md / .pdf / .docx / .pptx / .html / .csv 的文档加载与切分模块。
切块策略：段落优先 > 固定长度兜底，保留元数据（source/file_type/chunk_index/page）。
"""
import csv
import os
import re
from html.parser import HTMLParser
from typing import Optional

import pathspec

_IGNORE_FILE = ".doc_loader_ignore"
_IMAGE_EXTENSIONS = frozenset({".png", ".jpg", ".jpeg", ".webp"})
SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".pptx", ".html", ".htm", ".csv", ".xlsx"} | _IMAGE_EXTENSIONS


# ──────────────────────────────────────────────
# 文件读取
# ──────────────────────────────────────────────

def _read_txt_md(filepath: str) -> list[dict]:
    """读取 txt/md，返回单一 page（page=None）。"""
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    return [{"text": text, "page": None}]


def _table_to_text(headers: list[str], rows: list[list[str]]) -> str:
    """将表格格式化为 'header: value, ...' 样式的文本，跳过空值。"""
    formatted_rows: list[str] = []
    for row in rows:
        parts: list[str] = []
        for h, v in zip(headers, row):
            v_str = str(v).strip() if v else ""
            if v_str:
                parts.append(f"{h.strip()}: {v_str}")
        if parts:
            formatted_rows.append(", ".join(parts))
    return "\n".join(formatted_rows)


def _read_pdf(filepath: str) -> list[dict]:
    """读取 PDF，每页返回一个 page 字典，保留页码。优先使用 PyMuPDF，回退到 pypdf。"""
    # Try PyMuPDF first
    try:
        import fitz
    except ImportError:
        pass
    else:
        doc = fitz.open(filepath)
        try:
            pages = []
            for page_index, page in enumerate(doc, start=1):
                text = (page.get_text() or "").strip()
                if text:
                    pages.append({"text": text, "page": page_index})
            return pages
        finally:
            doc.close()

    # Fall back to pypdf
    try:
        from pypdf import PdfReader
    except ImportError:
        raise ImportError("请安装 PyMuPDF 或 pypdf：pip install pymupdf 或 pip install pypdf")

    reader = PdfReader(filepath)
    pages = []
    for page_index, page in enumerate(reader.pages, start=1):
        text = page.extract_text(extraction_mode="layout") or ""
        text = text.strip()
        if text:
            pages.append({"text": text, "page": page_index})
    return pages


def _read_docx(filepath: str) -> list[dict]:
    """读取 docx，返回合并段落与表格内容的单一 page（page=None）。"""
    try:
        from docx import Document
    except ImportError:
        raise ImportError("请安装 python-docx：pip install python-docx")

    doc = Document(filepath)
    text_parts: list[str] = []

    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    if paragraphs:
        text_parts.append("\n\n".join(paragraphs))

    for table in doc.tables:
        rows_data: list[list[str]] = []
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            rows_data.append(cells)
        if not rows_data:
            continue
        headers = rows_data[0]
        data_rows = rows_data[1:]
        table_text = _table_to_text(headers, data_rows)
        if table_text.strip():
            text_parts.append(table_text)

    text = "\n\n".join(text_parts)
    return [{"text": text, "page": None}]


def _read_pptx(filepath: str) -> list[dict]:
    """读取 PPTX，每张幻灯片返回一个 page 字典，保留页码（1-based）。提取正常文本框架和表格单元格文本。跳过空幻灯片。"""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("请安装 python-pptx：pip install python-pptx")

    prs = Presentation(filepath)
    slides_data: list[dict] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        text_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        text_parts.append(para_text)
            if shape.has_table:
                table = shape.table
                rows_data: list[list[str]] = []
                for row in table.rows:
                    cells = [cell.text for cell in row.cells]
                    rows_data.append(cells)
                if rows_data:
                    headers = rows_data[0]
                    data_rows = rows_data[1:]
                    table_text = _table_to_text(headers, data_rows)
                    if table_text.strip():
                        text_parts.append(table_text)
        text = "\n\n".join(text_parts)
        if text.strip():
            slides_data.append({"text": text, "page": slide_index})
    return slides_data


class _HTMLTextExtractor(HTMLParser):
    """从 HTML 中提取纯文本，忽略 script/style 标签。"""

    def __init__(self):
        super().__init__()
        self._text_parts: list[str] = []
        self._skip = False
        self._skip_tags = {"script", "style"}

    def handle_starttag(self, tag, attrs):
        if tag in self._skip_tags:
            self._skip = True

    def handle_endtag(self, tag):
        if tag in self._skip_tags:
            self._skip = False

    def handle_data(self, data):
        if not self._skip:
            stripped = data.strip()
            if stripped:
                self._text_parts.append(stripped)

    def get_text(self) -> str:
        return "\n\n".join(self._text_parts)


def _read_html(filepath: str) -> list[dict]:
    """读取 HTML，去除标签后返回纯文本。"""
    with open(filepath, "r", encoding="utf-8") as f:
        html_content = f.read()
    parser = _HTMLTextExtractor()
    parser.feed(html_content)
    text = parser.get_text()
    return [{"text": text, "page": None}] if text.strip() else []


def _read_csv(filepath: str, chunk_size: Optional[int] = None) -> list[dict]:
    """
    读取 CSV，每行输出 'key: value, key: value' 格式（含 header 作为字段名）。
    比纯逗号拼接对 LLM 语义更友好。

    当提供 chunk_size 且 >0 时，按完整行分组为多个 page，每组文本长度尽量 ≤ chunk_size
    （除非某单行本身超过 chunk_size）。
    """
    with open(filepath, "r", encoding="utf-8", newline="") as f:
        reader = csv.DictReader(f)
        rows = []
        for row in reader:
            parts = [f"{k.strip()}: {v.strip()}" for k, v in row.items() if v and v.strip()]
            if parts:
                rows.append(", ".join(parts))

    if not rows:
        return []

    if chunk_size is None or chunk_size <= 0:
        text = "\n\n".join(rows)
        return [{"text": text, "page": None}] if text.strip() else []

    # 按 chunk_size 按整行分组
    pages: list[dict] = []
    current_rows: list[str] = []
    current_len: int = 0
    sep = "\n\n"
    sep_len = len(sep)

    for row_text in rows:
        row_len = len(row_text)
        add_sep = sep_len if current_rows else 0

        if current_rows and current_len + add_sep + row_len > chunk_size:
            pages.append({"text": sep.join(current_rows), "page": None})
            current_rows = []
            current_len = 0
            add_sep = 0

        current_rows.append(row_text)
        current_len += add_sep + row_len

    if current_rows:
        pages.append({"text": sep.join(current_rows), "page": None})

    return pages


def _read_xlsx(filepath: str) -> list[dict]:
    """读取 XLSX，每个 sheet 返回一个 page 字典，页码为 sheet 序号（1-based）。

    使用 openpyxl（read_only / data_only），将每行格式化为 'header: value' 样式。
    """
    try:
        import openpyxl
    except ImportError:
        raise ImportError("请安装 openpyxl：pip install openpyxl")

    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    try:
        pages: list[dict] = []
        for sheet_index, sheet_name in enumerate(wb.sheetnames, start=1):
            ws = wb[sheet_name]
            rows_data: list[list[str]] = []
            for row in ws.iter_rows(values_only=True):
                row_values = [str(c) if c is not None else "" for c in row]
                rows_data.append(row_values)

            # 跳过完全空白的 sheet
            non_empty_rows = [r for r in rows_data if any(v.strip() for v in r)]
            if not non_empty_rows:
                continue

            # 第一行作为 header
            headers = non_empty_rows[0]
            data_rows = non_empty_rows[1:]

            if not data_rows:
                continue

            sheet_text_parts = [f"Sheet: {sheet_name}"]
            formatted = _table_to_text(headers, data_rows)
            if formatted.strip():
                sheet_text_parts.append(formatted)

            text = "\n\n".join(sheet_text_parts)
            if text.strip():
                pages.append({"text": text, "page": sheet_index})

        return pages
    finally:
        wb.close()


# ──────────────────────────────────────────────
# 图片提取（可选，用于 image caption）
# ──────────────────────────────────────────────


def _extract_pdf_images(filepath: str) -> list[tuple[bytes, int]]:
    """Extract embedded images from a PDF using PyMuPDF.

    Returns list of **(image_bytes, page_number)** tuples.
    Page numbers are 1-based.  Returns an empty list when PyMuPDF is
    unavailable or no images are found.
    """
    try:
        import fitz
    except ImportError:
        return []

    images: list[tuple[bytes, int]] = []
    doc = fitz.open(filepath)
    try:
        for page_index, page in enumerate(doc, start=1):
            for img_info in page.get_images(full=True):
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    images.append((image_bytes, page_index))
                except Exception:
                    continue
    finally:
        doc.close()
    return images


def _extract_docx_images(filepath: str) -> list[tuple[bytes, None]]:
    """Best-effort image extraction from DOCX via package relationships.

    Returns list of **(image_bytes, None)** tuples.  Page-level information
    is not available for DOCX files.
    """
    try:
        from docx import Document
    except ImportError:
        return []

    images: list[tuple[bytes, None]] = []
    try:
        doc = Document(filepath)
        for rel in doc.part.rels.values():
            reltype = getattr(rel, "reltype", "")
            if "image" in (reltype or ""):
                try:
                    blob = rel.target_part.blob
                    images.append((blob, None))
                except Exception:
                    continue
    except Exception:
        return []
    return images


def read_file_pages(filepath: str, chunk_size: Optional[int] = None) -> list[dict]:
    """
    读取文件，返回 list of {"text": str, "page": int|None}。
    page 对 PDF/PPTX 有效（1-based），其他格式为 None。
    chunk_size 仅用于 CSV 格式的按整行分组。
    """
    ext = os.path.splitext(filepath)[1].lower()
    if ext in {".txt", ".md"}:
        return _read_txt_md(filepath)
    if ext == ".pdf":
        return _read_pdf(filepath)
    if ext == ".docx":
        return _read_docx(filepath)
    if ext == ".pptx":
        return _read_pptx(filepath)
    if ext in {".html", ".htm"}:
        return _read_html(filepath)
    if ext == ".csv":
        return _read_csv(filepath, chunk_size=chunk_size)
    if ext == ".xlsx":
        return _read_xlsx(filepath)
    return []


# ──────────────────────────────────────────────
# Markdown 标题路径提取
# ──────────────────────────────────────────────

def _extract_heading_path(text: str) -> Optional[str]:
    """
    从 Markdown 文本中提取标题层级，返回如 "Chapter1 > Section2 > SubSection" 的字符串，
    如果没有标题则返回 None。
    """
    headings = re.findall(r"^(#{1,6})\s+(.+)", text, re.MULTILINE)
    if not headings:
        return None
    parts = [h[1].strip() for h in headings]
    return " > ".join(parts[:3])  # 最多3层


def _split_markdown_by_headings(text: str, chunk_size: int, chunk_overlap: int) -> list[tuple[Optional[str], str, int]]:
    """
    按 # 标题切分 Markdown，返回 [(heading_path, chunk_text, char_start), ...]。
    char_start 是 chunk_text 在 text 中的起始字符位置（不含 heading prefix）。
    维护一个标题栈，每遇到新标题就更新路径；同一标题节内部再按段落切。
    """
    heading_pattern = re.compile(r"^(#{1,6})\s+(.+)", re.MULTILINE)
    matches = list(heading_pattern.finditer(text))

    if not matches:
        # 没有标题，整篇按段落切
        chunks_with_pos = split_text_by_paragraphs(text, chunk_size, chunk_overlap)
        return [(None, c, pos) for c, pos in chunks_with_pos]

    # 按标题层级维护栈：[(level, title), ...]
    heading_stack: list[tuple[int, str]] = []
    sections: list[tuple[str, str, int]] = []  # (heading_path, section_text, section_offset)

    # 标题之前的前言部分
    preamble = text[:matches[0].start()].strip()
    if preamble:
        sections.append(("", preamble, 0))

    for i, m in enumerate(matches):
        level = len(m.group(1))  # # = 1, ## = 2, ...
        title = m.group(2).strip()

        # 弹出同级或更低级别的标题
        while heading_stack and heading_stack[-1][0] >= level:
            heading_stack.pop()
        heading_stack.append((level, title))

        # 节的文本范围：从当前标题到下一个同级或更高级标题
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        section_text = text[start:end].strip()

        if section_text:
            heading_path = " > ".join(t for _, t in heading_stack)
            sections.append((heading_path, section_text, start))

    # 每节内部按段落切块
    result: list[tuple[Optional[str], str, int]] = []
    for heading_path, section_text, section_offset in sections:
        chunks_with_pos = split_text_by_paragraphs(section_text, chunk_size, chunk_overlap)
        for chunk_text, chunk_pos in chunks_with_pos:
            result.append((heading_path if heading_path else None, chunk_text, section_offset + chunk_pos))

    return result


# ──────────────────────────────────────────────
# 切块函数
# ──────────────────────────────────────────────

def split_text_by_paragraphs(
    text: str, chunk_size: int, chunk_overlap: int
) -> list[tuple[str, int]]:
    """
    段落优先切块：先按空行分段，段落太长再按固定长度切。
    返回 [(chunk_text, char_start), ...]，char_start 是该块在 text 中的起始字符位置。
    """
    # 防御：overlap 必须严格小于 chunk_size，否则固定长度兜底时会死循环
    safe_overlap = min(chunk_overlap, chunk_size - 1) if chunk_size > 1 else 0

    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]

    results = []
    current = ""
    current_start = 0
    cursor = 0  # tracks our forward progress through text for position lookups

    for paragraph in paragraphs:
        # 使用 cursor 追踪原文进度，避免重复段落错误地匹配到更早的位置
        para_pos = text.find(paragraph, cursor)
        if para_pos < 0:
            para_pos = cursor
        cursor = para_pos + len(paragraph)

        if len(current) + len(paragraph) + 2 <= chunk_size:
            if not current:
                current_start = para_pos
            current = f"{current}\n\n{paragraph}" if current else paragraph
        else:
            if current:
                results.append((current, current_start))
            if len(paragraph) <= chunk_size:
                current = paragraph
                current_start = para_pos
            else:
                # 段落过长，固定长度兜底，使用 safe_overlap 保证 start 单调递增
                start = 0
                while start < len(paragraph):
                    end = start + chunk_size
                    results.append((paragraph[start:end], para_pos + start))
                    if end >= len(paragraph):
                        break
                    start = end - safe_overlap
                current = ""

    if current:
        results.append((current, current_start))

    return results


# ──────────────────────────────────────────────
# ignore 文件处理
# ──────────────────────────────────────────────

def _collect_ignore_specs(docs_dir: str) -> list[tuple[str, pathspec.PathSpec]]:
    specs: list[tuple[str, pathspec.PathSpec]] = []
    for root, _, files in os.walk(docs_dir):
        if _IGNORE_FILE in files:
            filepath = os.path.join(root, _IGNORE_FILE)
            with open(filepath, "r", encoding="utf-8") as f:
                spec = pathspec.PathSpec.from_lines("gitwildmatch", f)
                specs.append((root, spec))
    return specs


def _is_ignored(filepath: str, specs: list[tuple[str, pathspec.PathSpec]]) -> bool:
    for base, spec in specs:
        if filepath.startswith(base):
            rel = os.path.relpath(filepath, base)
            if spec.match_file(rel):
                return True
    return False


def load_document_images(docs_dir: str) -> list[dict]:
    """Extract image records from supported document files.

    Returned records are intentionally separate from text chunks because CLIP
    image embeddings must live in their own vector collection.
    """
    ignore_specs = _collect_ignore_specs(docs_dir)
    records: list[dict] = []

    for root, _, files in os.walk(docs_dir):
        for filename in sorted(files):
            ext = os.path.splitext(filename)[1].lower()
            if ext not in {".pdf", ".docx"}:
                continue

            filepath = os.path.join(root, filename)
            if _is_ignored(filepath, ignore_specs):
                continue

            relative_path = os.path.relpath(filepath, docs_dir)
            try:
                image_entries = (
                    _extract_pdf_images(filepath)
                    if ext == ".pdf"
                    else _extract_docx_images(filepath)
                )
            except Exception as e:
                print(f"[doc_loader] Image extraction failed for {relative_path}: {e}")
                continue

            for image_index, (image_bytes, page_num) in enumerate(image_entries):
                label = f"Image from {relative_path}"
                if page_num is not None:
                    label += f", page {page_num}"
                records.append(
                    {
                        "image_bytes": image_bytes,
                        "text": label,
                        "metadata": {
                            "source": relative_path,
                            "file_type": "clip_image",
                            "page": page_num,
                            "image_index": image_index,
                        },
                    }
                )

    return records


# ──────────────────────────────────────────────
# 主入口
# ──────────────────────────────────────────────

def load_documents(
    docs_dir: str,
    chunk_size: int = 800,
    chunk_overlap: int = 120,
    show_progress: bool = False,
    image_captioner=None,
    docs_lang: str = "zh",
    only_sources: Optional[set[str]] = None,
    parent_chunk_size: int | None = None,
) -> list[dict]:
    """
    遍历 docs_dir，加载所有支持格式的文档并切块。

    当提供 *image_captioner* 且文件为 PDF/DOCX 时，会将文档内嵌图片
    转为额外的 caption 文本块（metadata 中 file_type='image_caption'）。

    参数：
        only_sources: 若提供，仅加载相对路径在此集合中的文件。

    返回：
        list of {
            "text": str,
            "metadata": {
                "source": str,       # 相对路径
                "file_type": str,    # 扩展名（不含点）
                "chunk_index": int,  # 该文件内的 chunk 序号（0-based）
                "page": int|None,    # PDF/PPTX 页码/幻灯片号（1-based），其他为 None
            }
        }
    """
    ignore_specs = _collect_ignore_specs(docs_dir)
    # Heading prefix label is configurable per docs_lang to avoid awkward
    # Chinese-only (or English-only) labels in chunks.
    _heading_prefix_label = "标题路径：" if docs_lang == "zh" else "Heading: "
    all_chunks: list[dict] = []
    file_count = 0

    for root, _, files in os.walk(docs_dir):
        for filename in sorted(files):
            if filename == _IGNORE_FILE:
                continue

            ext = os.path.splitext(filename)[1].lower()
            if ext not in SUPPORTED_EXTENSIONS:
                continue

            filepath = os.path.join(root, filename)
            if _is_ignored(filepath, ignore_specs):
                continue

            relative_path = os.path.relpath(filepath, docs_dir)
            if only_sources is not None and relative_path not in only_sources:
                continue
            file_type = ext.lstrip(".")

            try:
                pages = read_file_pages(filepath, chunk_size=chunk_size)
            except Exception as e:
                print(f"[doc_loader] 跳过 {relative_path}：{e}")
                continue

            chunk_index = 0
            for page_data in pages:
                raw_text = page_data["text"]
                page_num = page_data["page"]

                if not raw_text.strip():
                    continue

                if ext == ".md":
                    # Markdown：按标题节切分，每个 chunk 带自己所属的标题路径
                    md_chunks = _split_markdown_by_headings(raw_text, chunk_size, chunk_overlap)
                    for heading_path, chunk_text, char_start in md_chunks:
                        if not chunk_text.strip():
                            continue
                        full_text = f"{_heading_prefix_label}{heading_path}\n\n{chunk_text}" if heading_path else chunk_text
                        meta = {
                            "source": relative_path,
                            "file_type": file_type,
                            "chunk_index": chunk_index,
                            "page": page_num,
                        }
                        if parent_chunk_size is not None:
                            meta["char_start"] = char_start
                        all_chunks.append({"text": full_text, "metadata": meta})
                        chunk_index += 1
                else:
                    text_chunks = split_text_by_paragraphs(
                        raw_text, chunk_size, chunk_overlap
                    )

                    for chunk_text, char_start in text_chunks:
                        if not chunk_text.strip():
                            continue
                        meta = {
                            "source": relative_path,
                            "file_type": file_type,
                            "chunk_index": chunk_index,
                            "page": page_num,
                        }
                        if parent_chunk_size is not None:
                            meta["char_start"] = char_start
                        all_chunks.append({"text": chunk_text, "metadata": meta})
                        chunk_index += 1

            # ── image caption chunks (opt-in, PDF / DOCX only) ──
            if image_captioner is not None and ext in (".pdf", ".docx"):
                try:
                    if ext == ".pdf":
                        img_entries = _extract_pdf_images(filepath)
                    else:
                        img_entries = _extract_docx_images(filepath)

                    if img_entries:
                        for image_index, (img_bytes, page_num) in enumerate(img_entries):
                            try:
                                hint = f"Image from {relative_path}"
                                if page_num is not None:
                                    hint += f", page {page_num}"
                                caption_text = image_captioner.caption(
                                    img_bytes, context_hint=hint
                                )
                            except Exception as e:
                                print(
                                    f"[doc_loader] Skip image {image_index}"
                                    f" in {relative_path}: {e}"
                                )
                                continue

                            if caption_text and caption_text.strip():
                                all_chunks.append(
                                    {
                                        "text": caption_text,
                                        "metadata": {
                                            "source": relative_path,
                                            "file_type": "image_caption",
                                            "chunk_index": chunk_index,
                                            "page": page_num,
                                            "image_index": image_index,
                                        },
                                    }
                                )
                                chunk_index += 1
                except Exception as e:
                    print(
                        f"[doc_loader] Image extraction failed"
                        f" for {relative_path}: {e}"
                    )

            # ── standalone image file captioning (opt-in) ──
            if image_captioner is not None and ext in _IMAGE_EXTENSIONS:
                try:
                    with open(filepath, "rb") as f:
                        img_bytes = f.read()
                    caption_text = image_captioner.caption(
                        img_bytes, context_hint=f"Image from {relative_path}"
                    )
                except Exception as e:
                    print(
                        f"[doc_loader] Skip image {relative_path}: {e}"
                    )
                else:
                    if caption_text and caption_text.strip():
                        all_chunks.append(
                            {
                                "text": caption_text,
                                "metadata": {
                                    "source": relative_path,
                                    "file_type": "image_caption",
                                    "chunk_index": chunk_index,
                                    "page": None,
                                    "image_index": 0,
                                },
                            }
                        )
                        chunk_index += 1

            file_count += 1
            if show_progress:
                print(
                    f"\r  {len(all_chunks)} chunks from {file_count} files...",
                    end="",
                    flush=True,
                )

    # ── Parent-child chunking ──
    if parent_chunk_size is not None and all_chunks:
        import hashlib

        def _generate_parent_id(source: str, parent_index: int, page=None) -> str:
            raw = f"{source}::page={page}::parent::{parent_index}"
            return hashlib.md5(raw.encode("utf-8")).hexdigest()[:12]

        # Group child chunks by (source, page)
        child_groups: dict[tuple, list[dict]] = {}
        for chunk in all_chunks:
            meta = chunk["metadata"]
            key = (meta["source"], meta.get("page"))
            child_groups.setdefault(key, []).append(chunk)

        # Cache file pages to avoid re-reading
        _pages_cache: dict[str, list[dict]] = {}

        parent_chunks: list[dict] = []

        for (source, page), group in child_groups.items():
            # Get raw text for this (source, page)
            if source not in _pages_cache:
                filepath = os.path.join(docs_dir, source)
                try:
                    _pages_cache[source] = read_file_pages(filepath, chunk_size=chunk_size)
                except Exception:
                    _pages_cache[source] = []

            raw_text = ""
            for page_data in _pages_cache[source]:
                if page_data["page"] == page:
                    raw_text = page_data["text"]
                    break
            if not raw_text.strip():
                continue

            # Split raw text into parent-sized chunks
            parent_chunks_with_pos = split_text_by_paragraphs(
                raw_text, parent_chunk_size, chunk_overlap
            )

            # Assign parent_ids and build parent chunks
            parent_id_map: dict[int, str] = {}  # parent_index -> parent_id
            for parent_idx, (parent_text, parent_char_start) in enumerate(parent_chunks_with_pos):
                if not parent_text.strip():
                    continue
                pid = _generate_parent_id(source, parent_idx, page)
                parent_id_map[parent_idx] = pid
                parent_chunks.append({
                    "text": parent_text,
                    "metadata": {
                        "source": source,
                        "file_type": group[0]["metadata"]["file_type"],
                        "chunk_index": parent_idx,
                        "page": page,
                        "chunk_role": "parent",
                        "parent_id": pid,
                    },
                })

            # Link each child chunk to its parent
            for chunk in group:
                child_start = chunk["metadata"].get("char_start")
                if child_start is None:
                    continue
                for parent_idx, (parent_text, parent_char_start) in enumerate(parent_chunks_with_pos):
                    if not parent_text.strip():
                        continue
                    parent_end = parent_char_start + len(parent_text)
                    if parent_char_start <= child_start < parent_end:
                        pid = parent_id_map.get(parent_idx)
                        if pid:
                            chunk["metadata"]["chunk_role"] = "child"
                            chunk["metadata"]["parent_id"] = pid
                        break

        # Remove char_start from child metadata (internal use only)
        for chunk in all_chunks:
            chunk["metadata"].pop("char_start", None)

        all_chunks.extend(parent_chunks)

    return all_chunks
